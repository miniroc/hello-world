"""
https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
http://pbpython.com/creating-powerpoint.html
"""
from __future__ import print_function

import re
import urllib.request
import urllib.error
from urllib.request import ProxyHandler
from urllib.request import build_opener
# import argparse
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

from bs4 import BeautifulSoup


def download_index(url, file_path):
    feature_dicts_list = []
    try:
        #using header to simulate browser
#         headers = {}
#         headers['User-Agent'] = "Mozilla/5.0 (Windows NT 6.1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/67.0.3396.99 Safari/537.36"
        headers = {'User-Agent': "User-Agent: Mozilla/5.0 (Windows NT 6.1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/67.0.3396.99 Safari/537.36"}
#         req = urllib.request.Request("https://www.phonearena.com/phones/carriers/MetroPCS", headers=headers)

        req = urllib.request.Request(url, headers=headers)
        html = urllib.request.urlopen(req)
        result = html.read().decode('utf-8')
#         with open("ProductIndex.xml", 'w+') as f:
#             f.write(result)
#         print(result)

    except urllib.error.URLError as e:
        if hasattr(e, 'reason'):
            print("error reason: " + str(e.reason))
    except urllib.error.HTTPError as e:
        if hasattr(e, 'code'):
            print("error reason: " + str(e.code))
    else:
        print("urllib.request.urlopen() Succeeded")
    
    soup = BeautifulSoup(result,'html.parser')
    if(soup is None):
        print("soup failed")
        return
    
#   get product url
    url_div_tags = soup.findAll('div', class_='quicklookdiv')
    if(url_div_tags is None):
        print("hyper_link_tag is None")
        return
    for url_div_tag in url_div_tags:
#         print(url_div_tag)
        hyper_link_tags = url_div_tag.findAll('a', href=True)
        if(hyper_link_tags is None):
            print("hyper_link_tag is None")
            return
        # get link address from a Tag object
        for hyper_link_tag in hyper_link_tags:
            hyper_link = "https://www.phonearena.com" + hyper_link_tag['href']
            print(hyper_link)
#             start recursively open product link and extract spec
            feature_dict = downloadpages(hyper_link)
#             add obtained spec into an existing ppt file
            if(feature_dict is None):
                print("failed to get feature list")
                return
            feature_dicts_list.append(feature_dict)

    addslides(file_path, file_path, feature_dicts_list)

def downloadpages(url):
    '''html format only works for phonearena detailed spec page like https://www.phonearena.com/phones/LG-Stylo-4_id10876'''
    spec_dict = {}
    try:
        #using header to simulate browser
#         headers = {}
#         headers['User-Agent'] = "Mozilla/5.0 (Windows NT 6.1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/67.0.3396.99 Safari/537.36"
        headers = {'User-Agent': "User-Agent: Mozilla/5.0 (Windows NT 6.1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/67.0.3396.99 Safari/537.36"}
#         req = urllib.request.Request("https://www.phonearena.com/phones/carriers/MetroPCS", headers=headers)

        req = urllib.request.Request(url, headers=headers)
        html = urllib.request.urlopen(req)
        result = html.read().decode('utf-8')
#         with open("result.xml", 'w+') as f:
#             f.write(result)
#         print(result)

    except urllib.error.URLError as e:
        if hasattr(e, 'reason'):
            print("error reason: " + str(e.reason))
    except urllib.error.HTTPError as e:
        if hasattr(e, 'code'):
            print("error reason: " + str(e.code))
    else:
        print("urllib.request.urlopen() Succeeded")
    
    
# use re to find tags
# #     <strong class=" s_lv_1 ">Dimensions:</strong>
# #     spec_item = re.findall(r'<strong class=" s_lv_1 ">(.*?)</strong>', result) # add ? to denote non-greedy mode
#     spec_content = re.findall(r'<div class=" clear specs-holder">(.*?)</div>', result) # add ? to denote non-greedy mode
# #     print(len(spec_result))
# #     print(spec_result)
# #     print(spec_result[1])
#     with open("data.xml", 'w+') as f:
# #         for i in range(0, len(spec_item), 1):
# #             f.write(spec_item[i])
#         for i in range(0, len(spec_content), 1):
#             f.write(spec_content[i])
    
    # use beautifulsoup to organize html
    soup = BeautifulSoup(result,'html.parser')
    if(soup is None):
        print("soup failed")
        return
    
#   get product ID
    pic_div = soup.find('div', class_='quicklook')
    hyper_link_tag = pic_div.find('a', href=True)
    # get link address from a Tag object
    hyper_link = hyper_link_tag['href']
#     temp_file_name = hyper_link[-5:] # will cause conflic pic name
#     print(hyper_link_tag['href'])
    temp_file_name = hyper_link.split("/")[-1]
    urllib.request.urlretrieve(hyper_link, temp_file_name) 
    spec_dict['phone_id'] = temp_file_name
    
#   get product name
    title_tag = soup.find('title')
    spec_dict['phone_name'] = title_tag.text
    
#     get product release date
    release_date_tag = soup.find('div', class_='metainfo')
#     some product don't have Posted: or Release date, find either one is OK
    release_dates = re.findall(r'Posted:.*?20[0-9][0-9]|Release date:.*?20[0-9][0-9]',release_date_tag.text)
    if(release_dates is not None):
        release_date = re.sub("Posted:","",release_dates[0])
        release_date = re.sub("Release date:","",release_date)
#         print("release data: ", release_date)
        spec_dict['Release_date'] = release_date
#     release_date_tag.text.
    
    carriers_tag = soup.find("div", class_='carriers')
    if(carriers_tag is not None):
        carrier_tags = carriers_tag.findAll('a')
        spec_dict['Carriers'] = ""
        for carrier_tag in carrier_tags:
            spec_dict['Carriers'] = spec_dict['Carriers'] + carrier_tag.text + " "
    
#     for each in soup.find_all('a', class_=["s_block_4 s_block_4_s115  s_fst  clearfix","s_block_4 s_block_4_s115     clearfix"]):
#     for each in soup.find_all('div', class_=' clear specs-holder'):
#     with open("data.xml", 'w+') as f:
#         for each in soup.find_all('div', attrs={'class': ' clear specs-holder'}):
#             # Take out the <div> of name and get its value
#             name = each.text.strip() # strip() is used to remove starting and trailing
#             print(name)
#             f.write(name)

# <div class="s_specs_box s_box_4" id="design_cPoint">
# <h2 class="htitle">Design</h2>
# <ul>
# <li class="s_lv_1 field-6">
#     <strong class=" s_lv_1 ">Dimensions:</strong>
#     <ul class=" s_lv_1 ">
#         <li><span title='Big Dimensions' class="s_size_rating s_size_rating_s1 s_mr_5"><span></span></span>
#             6.30 x 3.06 x 0.32 inches  (160 x 77.7 x 8.1 mm)
#         </li>
#     </ul>
# </li>
# <li class="s_lv_1 field-166">
#     <strong class=" s_lv_1 ">Weight:</strong>
#     <ul class=" s_lv_1 ">
#         <li>
#             <span title='Average Weight. The average weight is: 5.7 oz (161 g).' class="s_weight_rating s_weight_rating_s2 s_mr_5"><span></span></span>
#             6.07 oz  (172 g)<br />
#             <span class="s_f_11 gray_9">the <span class="s_tooltip_info blue_1 tooltip" title="<p>Average represents the mean value, calculated from all phones, from the last year.">average</span> is 5.7 oz (161 g)</span>
#         </li>
#     </ul>
# </li>

# get detailed spec
    spec_fields = soup.findAll('div', attrs={'class':'s_specs_box s_box_4'})
    if(spec_fields is None):
        print("field is none")
        return
#     print(spec_fields)
    
    for field in spec_fields:
        strongs = field.findAll('strong')
        for strong in strongs:
#         <span class='s_tooltip_anchor'>Capacity:</span>
            tooltip = strong.find('span', class_='s_tooltip_anchor')
            if(tooltip is not None):
#                 print("span: "+tooltip.text)
                value = strong.findNextSibling().find('li')
#                 print("value: "+value.text)
                
                spec_dict[tooltip.text.replace(":","")] = value.text.replace("\n","")
            else:        
#                 print("stong: "+strong.text)
                value = strong.findNextSibling().find('li')
                [x.extract() for x in value.findAll('span')]
#                 print("value: " + value.text)
                
                spec_dict[strong.text.replace(":","")] = value.text.replace("\n","")
#             print("\n")
#     for k, v in spec_dict.items():
#         print("({key}: {value})".format(key=k, value=v))
        
#         <ul class=" s_lv_1 ">
#             <li>3300 mAh</li>
#         </ul>
#         value = field.find('ul', class_=' s_lv_1 ')
#         print("ul: "+value.text)
#     display_size = soup.find('span', attrs={'title': 'Big display'}).parent
#     print(display_size.text)
#     [''LCD: Physical size:', 'Battery: Capacity', 'Size: Dimensions:', 'Chipset: System chip:', 'Bands', 'Camera', 'Connectivity', 'ROM/RAM',
#                    'Fingerprint', 'Waterproof', 'OS', 'CMF', 'Launch time', 'Carrier']
#    merge some of the fileds
# find if key exist
#     if('Resolution' in spec_dict.keys()):
# Another way to find if key exist
    if('Physical size' in spec_dict.keys()):
        if('Resolution' in spec_dict.keys()):
            spec_dict['Physical size'] = spec_dict['Physical size'] + ", " + spec_dict['Resolution'] 
            del spec_dict['Resolution']
#             rename dict key name: Physical size -> Display Size
            spec_dict['Screen size'] = spec_dict.pop('Physical size')
    if('Capacity' in spec_dict.keys()): 
#         rename dict key name: Physical size -> Display Size
        spec_dict['Battery'] = spec_dict.pop('Capacity') 

    if('Camera' in spec_dict.keys()):
        if('Aperture size' in spec_dict.keys()):
            spec_dict['Camera'] = spec_dict['Camera'] + "(" + spec_dict['Aperture size'] + ")"
            del spec_dict['Aperture size']
            
            if('Front-facing camera' in spec_dict.keys()):
                spec_dict['Camera'] = spec_dict['Camera'] + " + " + spec_dict['Front-facing camera']
                del spec_dict['Front-facing camera']
                 
    if('System memory' in spec_dict.keys()):
        if('Built-in storage' in spec_dict.keys()):
            spec_dict['System memory'] = spec_dict['System memory'] + ", " + spec_dict['Built-in storage'] 
            del spec_dict['Built-in storage']
#             rename dict key name
            spec_dict['Memory'] = spec_dict.pop('System memory')        
    
    if('GSM' in spec_dict.keys()):
        if('UMTS' in spec_dict.keys()):
            spec_dict['GSM'] = "GSM: " + spec_dict['GSM'] + "\nUMTS:" + spec_dict['UMTS']
            del spec_dict['UMTS']
            if('LTE (FDD)' in spec_dict.keys()):
                spec_dict['GSM'] = spec_dict['GSM'] + "\nLTE: " + spec_dict['LTE (FDD)']
                del spec_dict['LTE (FDD)']
                if('LTE (TDD)' in spec_dict.keys()):
                    spec_dict['GSM'] = spec_dict['GSM'] +" " + spec_dict['LTE (TDD)']
                    del spec_dict['LTE (TDD)']
            spec_dict['Bands'] = spec_dict.pop('GSM')
       
    return spec_dict

def addslides(inputfile, outputfile, feature_dicts_list):
    """ read a pptx and add slides
    """
    # still missing launch date, price, carriers. , , is not existed in every page
    keys = ["Screen size","Battery","Dimensions","Screen-to-body ratio","System chip","Camera",
            "Memory","Biometrics","Bluetooth","Wi-Fi", "Bands","OS","Materials","Release_date", 'Carriers']
    
    prs = Presentation(inputfile)
    # slide_layouts is a slide, Presentation.slides, slide.shapes, shapes.title, shapes.placeholders[]
    spec_slide_layout = prs.slide_layouts[2] # 2 is no.2 mother template
    
#     add slides according to feature_dicts_list numbers
    for spec_dict in feature_dicts_list:
        slide = prs.slides.add_slide(spec_slide_layout)
    
    # add product name    
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = spec_dict['phone_name']
        
    # shapes.placeholder -- body
#     body_shape = shapes.placeholders[12]
#     tf = body_shape.text_frame
#     tf.text = 'Presentation.slides.shapes.placeholders[].text_frame'
#      
#     p = tf.add_paragraph()
#     p.text = 'Presentation.slides.shapes.placeholders[].text_frame.paragraph level 1'
#     p.level = 1
#      
#     p = tf.add_paragraph()
#     p.text = 'Presentation.slides.shapes.placeholders[].text_frame.paragraph level 2'
#     p.level = 2
    
    # add product ID
        img_path = spec_dict['phone_id']    
        left = Inches(0.01)
        top = Inches(0.6)
        height = Inches(5)
        slide.shapes.add_picture(img_path, left, top, height=height)
        
    # add spec table
        rows = len(keys)+1
        cols = 2
        left = Inches(4.0)
        top = Inches(0.2)
        width = Inches(2.0)
        height = Inches(0.1)
         
        table = shapes.add_table(rows, cols, left, top, width, height).table
        
        # set column widths
        table.columns[0].width = Inches(1.6)
        table.columns[1].width = Inches(4.2)
    
        # write column headings
        table.cell(0, 0).text = 'Item'
        table.cell(0, 1).text = 'Spec'
    
        i=1
        for key in keys:
            # check if this items exist
            if(key in spec_dict.keys()):        
                key_frame = table.cell(i,0).text_frame
                key_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                p = key_frame.paragraphs[0]
                run = p.add_run()
                run.text = key    
                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)
                font.bold = False
                font.italic = None  # cause value to be inherited from theme
        #       
                value_frame = table.cell(i,1).text_frame
                value_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                p = value_frame.paragraphs[0]
                run = p.add_run()
                run.text = spec_dict[key]
                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)
                font.bold = False
                font.italic = None  # cause value to be inherited from theme
                
                i += 1        
        
# Store all the items
#     i=1
#     for k, v in spec_dict.items():
#         key_frame = table.cell(i,0).text_frame
#         key_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#         p = key_frame.paragraphs[0]
#         run = p.add_run()
#         run.text = k    
#         font = run.font
#         font.name = 'Calibri'
#         font.size = Pt(12)
#         font.bold = False
#         font.italic = None  # cause value to be inherited from theme
# #       
#         value_frame = table.cell(i,1).text_frame
#         value_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#         p = value_frame.paragraphs[0]
#         run = p.add_run()
#         run.text = v
#         font = run.font
#         font.name = 'Calibri'
#         font.size = Pt(12)
#         font.bold = False
#         font.italic = None  # cause value to be inherited from theme
#         
#         i += 1
  
    prs.save(outputfile)

if __name__ == "__main__":
    inputPPT = "C:\\1.Study Video\\eclipse_workspace\\hello-world\\WebCraw2PPT\\input.pptx"
#     feature_dict = downloadpages("https://www.phonearena.com/phones/LG-Aristo-2_id10782")
#     feature_dict_list = [feature_dict]
#     addslides(inputPPT,inputPPT, feature_dict_list)
    
    download_index("https://www.phonearena.com/phones/carriers/MetroPCS",inputPPT)

    
    